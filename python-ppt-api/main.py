"""
Layah.ai – Python PPT Generation API
======================================
Uses python-pptx to open an uploaded school .pptx template, clear its existing
content slides (keeping the slide master + layouts intact), then inject 13
lesson slides with AI-generated content while preserving the school's visual
design (background, colours, logo, fonts).

Endpoints
---------
GET  /health          – liveness check
POST /generate-ppt    – multipart: field "template" (.pptx file) + field "slides" (JSON)
"""

import copy
import io
import json
import os
import tempfile
import traceback

from flask import Flask, jsonify, request, send_file
from flask_cors import CORS

from lxml import etree

from pptx import Presentation
from pptx.oxml.ns import qn, nsmap
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

# ── XML namespaces ────────────────────────────────────────────────────────────
PML_NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAW_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ── Slide type map ────────────────────────────────────────────────────────────
SLIDE_TYPES = {
    0:  "title",
    1:  "starter",
    2:  "objectives",
    3:  "vocabulary",
    4:  "context",
    5:  "main",
    6:  "main2",
    7:  "activity",
    8:  "plenary",
    9:  "differentiated",
    10: "exit_ticket",
    11: "success_criteria",
    12: "closing",
}

TITLE_SLIDE_INDICES = {0, 12}


# ── Background helpers ────────────────────────────────────────────────────────

def get_slide_bg_element(slide):
    """
    Return the <p:bg> lxml element from a slide's <p:cSld>, or None.
    This element carries the slide-level background (colour fill, image, etc.)
    that designers often put directly on each slide rather than in the master.
    """
    cSld = slide._element.find(f"{{{PML_NS}}}cSld")
    if cSld is None:
        return None
    return cSld.find(f"{{{PML_NS}}}bg")


def apply_bg_to_slide(target_slide, bg_element):
    """
    Deep-copy bg_element and insert it as the first child of the target
    slide's <p:cSld>.  Replaces any existing <p:bg>.
    """
    if bg_element is None:
        return
    cSld = target_slide._element.find(f"{{{PML_NS}}}cSld")
    if cSld is None:
        return
    # Remove any existing bg
    old_bg = cSld.find(f"{{{PML_NS}}}bg")
    if old_bg is not None:
        cSld.remove(old_bg)
    # Insert the copied bg as the first element (before spTree)
    cSld.insert(0, copy.deepcopy(bg_element))


def get_master_bg_element(prs):
    """Return the <p:bg> element from the first slide master, if any."""
    try:
        master = prs.slide_masters[0]
        cSld = master._element.find(f"{{{PML_NS}}}cSld")
        if cSld is not None:
            return cSld.find(f"{{{PML_NS}}}bg")
    except Exception:
        pass
    return None


def collect_best_bg(prs, template_slides_before_removal):
    """
    Find the best background element to stamp onto every new slide.
    Priority: first template slide bg → slide master bg → None
    """
    # Try first content slide
    for s in template_slides_before_removal:
        bg = get_slide_bg_element(s)
        if bg is not None:
            print("[ppt-api] Using background from first template slide")
            return bg
    # Try slide master
    bg = get_master_bg_element(prs)
    if bg is not None:
        print("[ppt-api] Using background from slide master")
        return bg
    print("[ppt-api] No explicit background found in template (will rely on layout inheritance)")
    return None


# ── Slide removal ─────────────────────────────────────────────────────────────

def remove_all_slides(prs):
    """
    Remove every content slide from a presentation while keeping the
    slide masters and layouts (which carry the template's visual design).
    Returns the list of slide objects that existed before removal so their
    backgrounds can be copied to new slides.
    """
    sld_id_lst = prs.slides._sldIdLst
    existing   = list(sld_id_lst)
    slides_snapshot = list(prs.slides)   # snapshot for bg extraction

    print(f"[ppt-api] Removing {len(existing)} existing slide(s) from template...")
    for sld_id in existing:
        r_id = sld_id.get(qn("r:id"))
        if r_id:
            try:
                prs.part.drop_rel(r_id)
                print(f"  └─ Dropped slide rel: {r_id}")
            except Exception as e:
                print(f"  └─ Warning: could not drop rel {r_id}: {e}")
        sld_id_lst.remove(sld_id)

    print(f"[ppt-api] Slides after removal: {len(prs.slides)}")
    return slides_snapshot


# ── Layout selection ──────────────────────────────────────────────────────────

def find_layout(prs, name_hints: list):
    """
    Return the first layout whose name contains any hint (case-insensitive).
    Fallback: index 1 (usually 'Title and Content'), then 0.
    """
    for layout in prs.slide_layouts:
        ln = layout.name.lower()
        for hint in name_hints:
            if hint.lower() in ln:
                return layout
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


# ── Text helpers ──────────────────────────────────────────────────────────────

def set_placeholder_text(ph, lines: list, font_size_pt=None, bold: bool = False) -> None:
    """Clear a placeholder and fill it line by line with proper alignment."""
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        if font_size_pt:
            run.font.size = Pt(font_size_pt)
        if bold:
            run.font.bold = True


def add_text_box_safe(slide, text: str, left, top, width, height,
                      font_size_pt: int = 18, bold: bool = False) -> None:
    """
    Add a manual text box with safe margins and proper alignment.
    Splits text on newlines; each line becomes its own paragraph.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    lines = [l for l in text.split("\n") if l.strip()]
    if not lines:
        lines = [text.strip()]

    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold


# ── Slide builder ─────────────────────────────────────────────────────────────

def build_slide(prs, layout, slide_idx: int, slide_data: dict,
                slide_w, slide_h, bg_element=None) -> None:
    """
    Add one lesson slide using the template layout.
    The new slide inherits the full master/layout design automatically.
    If an explicit bg_element was found in the template, it is also stamped on.
    """
    title_text   = slide_data.get("title",   "").strip()
    content_text = slide_data.get("content", "").strip()
    slide_type   = SLIDE_TYPES.get(slide_idx, "content")

    print(f"[ppt-api] Slide {slide_idx + 1:02d} ({slide_type}): {title_text[:50]!r}")

    # ── Add the slide — inherits full visual design from layout → master ──────
    slide = prs.slides.add_slide(layout)

    # ── Stamp background (in case it was slide-level, not master-level) ───────
    if bg_element is not None:
        apply_bg_to_slide(slide, bg_element)

    # ── Safe layout margins ───────────────────────────────────────────────────
    margin   = Inches(0.5)
    title_h  = Inches(1.1)
    body_top = Inches(1.75)
    body_h   = slide_h - body_top - Inches(0.5)   # leave 0.5" bottom margin
    body_w   = slide_w - 2 * margin

    placed_title   = False
    placed_content = False

    # ── Fill placeholders first (best quality — inherits layout fonts/colour) ─
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        print(f"  └─ placeholder idx={idx} name={ph.name!r}")

        if idx == 0 and not placed_title and title_text:
            set_placeholder_text(ph, [title_text], bold=True)
            placed_title = True

        elif idx in (1, 2, 13, 14, 15) and not placed_content and content_text:
            lines = [l for l in content_text.split("\n") if l.strip()]
            set_placeholder_text(ph, lines)
            placed_content = True

    # ── Fallback text boxes when layout has no matching placeholder ───────────
    if not placed_title and title_text:
        print("  └─ No title placeholder — adding fallback text box")
        add_text_box_safe(slide, title_text,
                          left=margin, top=margin,
                          width=body_w, height=title_h,
                          font_size_pt=32, bold=True)

    if not placed_content and content_text:
        print("  └─ No content placeholder — adding fallback text box")
        add_text_box_safe(slide, content_text,
                          left=margin, top=body_top,
                          width=body_w, height=body_h,
                          font_size_pt=18)

    # ── Speaker notes ─────────────────────────────────────────────────────────
    notes_text = slide_data.get("speakerNotes", "").strip()
    if notes_text:
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
        except Exception:
            pass


# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "message": "Layah PPT API is running"})


@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    # ── 1. Confirm template receipt ───────────────────────────────────────────
    template_file = request.files.get("template")
    slides_json   = request.form.get("slides")

    if template_file:
        template_bytes = template_file.read()
        print(f"[ppt-api] ✔ Template file received — filename={template_file.filename!r}, size={len(template_bytes)} bytes")
    else:
        print("[ppt-api] ✘ No template file received!")
        return jsonify({"error": "Missing 'template' file in multipart form."}), 400

    if not slides_json:
        print("[ppt-api] ✘ No 'slides' JSON field received!")
        return jsonify({"error": "Missing 'slides' JSON field in form data."}), 400

    if len(template_bytes) < 1000:
        print(f"[ppt-api] ✘ Template file is too small ({len(template_bytes)} bytes) — likely corrupt")
        return jsonify({"error": f"Template file appears corrupt ({len(template_bytes)} bytes)."}), 400

    try:
        payload = json.loads(slides_json)
    except json.JSONDecodeError as exc:
        return jsonify({"error": f"Invalid JSON in 'slides': {exc}"}), 400

    slides_list = payload.get("slides", [])
    if not slides_list:
        return jsonify({"error": "'slides' array is empty."}), 400

    topic   = payload.get("topic",       "Lesson")
    subject = payload.get("subject",     "")
    grade   = payload.get("grade",       "")
    teacher = payload.get("teacherName", "Teacher")

    print(f"[ppt-api] ══ generate-ppt: topic={topic!r}, subject={subject!r}, grade={grade!r}, slides={len(slides_list)} ══")

    tmp_path = None
    try:
        # ── 2. Save template to a real file (more reliable than BytesIO) ──────
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(template_bytes)
            tmp_path = tmp.name
        print(f"[ppt-api] Template saved to temp file: {tmp_path}")

        # ── 3. Open with python-pptx ──────────────────────────────────────────
        prs = Presentation(tmp_path)

        print(f"[ppt-api] ── Template info ──────────────────────────────────")
        print(f"  Existing slides  : {len(prs.slides)}")
        print(f"  Slide masters    : {len(prs.slide_masters)}")
        print(f"  Slide layouts    : {len(prs.slide_layouts)}")
        print(f"  Slide dimensions : {prs.slide_width} × {prs.slide_height} EMUs")

        if len(prs.slide_layouts) == 0:
            return jsonify({"error": "Template has no slide layouts."}), 500

        print(f"[ppt-api] ── Slide layouts ({len(prs.slide_layouts)}) ────────────────")
        for i, layout in enumerate(prs.slide_layouts):
            ph_count = len(list(layout.placeholders))
            print(f"  Layout [{i:02d}]: {layout.name!r}  (placeholders: {ph_count})")

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # ── 4. Capture background BEFORE removing slides ──────────────────────
        #     This is the critical step: we extract the slide-level background
        #     from the template's existing slides so we can stamp it on every
        #     new slide we create.  If no slide-level bg exists, the master/
        #     layout background will still be inherited automatically.
        bg_element = collect_best_bg(prs, list(prs.slides))

        # ── 5. Remove existing content slides (masters + layouts stay) ────────
        remove_all_slides(prs)

        # ── 6. Choose layouts from the template ───────────────────────────────
        title_layout   = find_layout(prs, ["title slide", "title,", "cover", "section"])
        content_layout = find_layout(prs, ["title and content", "title, content", "content"])

        print(f"[ppt-api] ✔ Title layout  : {title_layout.name!r}")
        print(f"[ppt-api] ✔ Content layout: {content_layout.name!r}")

        # ── 7. Generate 13 lesson slides ──────────────────────────────────────
        for idx, slide_data in enumerate(slides_list):
            slide_type = SLIDE_TYPES.get(idx, "content")
            layout     = title_layout if slide_type in ("title", "closing") else content_layout
            build_slide(prs, layout, idx, slide_data, slide_w, slide_h, bg_element)

        print(f"[ppt-api] Final slide count: {len(prs.slides)}")

        # ── 8. Serialise and return ───────────────────────────────────────────
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        size = output.getbuffer().nbytes
        print(f"[ppt-api] ══ PPT generated successfully: {size} bytes ══")

        return send_file(
            output,
            as_attachment=True,
            download_name="lesson-presentation.pptx",
            mimetype=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        )

    except Exception as exc:
        print(f"[ppt-api] ✘ ERROR: {exc}")
        traceback.print_exc()
        return jsonify({"error": str(exc)}), 500

    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
                print(f"[ppt-api] Temp file cleaned up: {tmp_path}")
            except Exception:
                pass


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    port  = int(os.environ.get("PORT", 5001))
    debug = os.environ.get("FLASK_DEBUG", "0") == "1"
    print(f"[ppt-api] Starting on port {port} (debug={debug})")
    app.run(host="0.0.0.0", port=port, debug=debug)
